"""
Generalized document parsing and taxonomy mapping service.

Handles PDF, PPTX, XLSX, and CSV files — detects the firm's reporting structure
and maps their categories to SentioCap's L1-L4 taxonomy.

Works with any FP&A format: NEVER assumes a specific structure.
"""
from __future__ import annotations

import io
import json
import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Optional

import anthropic
import pandas as pd

from config import settings

logger = logging.getLogger(__name__)

MODEL = "claude-3-5-sonnet-20241022"

# ---------------------------------------------------------------------------
# Dataclasses
# ---------------------------------------------------------------------------


@dataclass
class ExtractedTable:
    headers: list[str]
    rows: list[list[Any]]
    page_or_slide: int
    title: str = ""

    def to_dict(self) -> dict:
        return {
            "title": self.title,
            "page_or_slide": self.page_or_slide,
            "headers": self.headers,
            "row_count": len(self.rows),
            "sample_rows": self.rows[:5],
        }


@dataclass
class Section:
    heading: str
    content: str
    level: int = 1  # heading depth


@dataclass
class ParsedDocument:
    raw_text: str
    tables: list[ExtractedTable]
    sections: list[Section]
    metadata: dict


@dataclass
class HierarchyLevel:
    name: str          # e.g. "Division", "Department", "GL Account"
    depth: int         # 1 = top level
    example_values: list[str] = field(default_factory=list)


@dataclass
class DetectedCategory:
    name: str
    level: str           # which hierarchy level it belongs to
    parent: Optional[str]
    example_items: list[str] = field(default_factory=list)
    total_amount: Optional[float] = None


@dataclass
class DetectedStructure:
    structure_tier: int          # 1=highly structured, 2=semi-structured, 3=unstructured
    hierarchy_levels: list[HierarchyLevel]
    has_rtb_ctb_split: bool
    has_department_breakdown: bool
    has_investment_tracking: bool
    has_kpis: bool
    detected_categories: list[DetectedCategory]
    currency: str
    fiscal_period: str
    confidence: float


@dataclass
class CategoryMapping:
    source_name: str
    source_level: str
    mapped_l1: str
    mapped_l2: str
    mapped_l3: str
    mapped_l4: str
    confidence: float
    reasoning: str
    needs_review: bool
    allocation_pct: float = 1.0    # if a source spans multiple targets


@dataclass
class TaxonomyMapping:
    mappings: list[CategoryMapping]
    unmapped: list[str]
    summary: str


@dataclass
class ClassifiedLineItem:
    source_description: str
    source_category: str
    amount: float
    period: str
    classified_l1: str
    classified_l2: str
    classified_l3: str
    classified_l4: str
    confidence: float


@dataclass
class ClassifiedData:
    line_items: list[ClassifiedLineItem]
    totals: dict
    insights: list[str]


# ---------------------------------------------------------------------------
# Taxonomy reference (used in prompts)
# ---------------------------------------------------------------------------

TAXONOMY_REFERENCE = """
## SentioCap Investment Intent Taxonomy

### L1 — Intent
- RTB (Run The Business): Operational spend maintaining existing capabilities — software maintenance, IT support, recurring SaaS subscriptions, infrastructure hosting for current systems.
- CTB (Change The Business): Investment spend building new capabilities — new software development, digital transformation, new platforms, AI/ML initiatives, R&D. Expenses capitalized under ASC 350-40 must be CTB.

### L2 — Category
- L2-INF: Infrastructure & Platforms — servers, storage, cloud compute, networking hardware, hosting.
- L2-APP: Applications & Software — business applications, SaaS, ERP/CRM, software licenses.
- L2-DAT: Data & Analytics — data warehousing, BI tools, analytics platforms, data engineering.
- L2-SEC: Security & Compliance — cybersecurity, compliance programs, identity management, pen testing.
- L2-OPS: IT Operations & Support — help desk, ITSM, monitoring, incident response.
- L2-PEO: People & Talent — IT staff salaries, contractors, consulting, training, certifications.
- L2-INO: Innovation & R&D — research, proof-of-concepts, emerging tech, AI/ML experimentation.
- L2-PRO: Process & Transformation — process redesign, change management, digital transformation programs.

### L3 — Domain
- L3-CLD: Cloud & Hosting
- L3-NET: Networking & Connectivity
- L3-SFT: Software Licensing
- L3-DEV: Development & Engineering
- L3-ANL: Analytics & BI
- L3-HCM: Human Capital Management
- L3-CYB: Cybersecurity
- L3-CHG: Change Management & Training

### L4 — Activity
- L4-CLD-001: Public Cloud Compute (L3-CLD)
- L4-CLD-002: Cloud Storage (L3-CLD)
- L4-CLD-003: Cloud Networking (L3-CLD)
- L4-NET-001: WAN / SD-WAN (L3-NET)
- L4-NET-002: LAN Infrastructure (L3-NET)
- L4-SFT-001: SaaS Subscriptions (L3-SFT)
- L4-SFT-002: On-Prem Software Licenses (L3-SFT)
- L4-DEV-001: Application Development (L3-DEV)
- L4-DEV-002: DevOps Tooling (L3-DEV)
- L4-ANL-001: BI & Reporting Tools (L3-ANL)
- L4-ANL-002: Data Engineering (L3-ANL)
- L4-HCM-001: IT Staff Salaries (L3-HCM)
- L4-HCM-002: Contractor / Consulting (L3-HCM)
- L4-HCM-003: Training & Certification (L3-HCM)
- L4-CYB-001: Security Software (L3-CYB)
- L4-CYB-002: Penetration Testing (L3-CYB)
- L4-CHG-001: Change Management (L3-CHG)
- L4-CHG-002: Process Redesign (L3-CHG)
"""

# ---------------------------------------------------------------------------
# AI Prompts
# ---------------------------------------------------------------------------

STRUCTURE_DETECTION_SYSTEM = """You are analyzing a financial document to detect its reporting structure.
Companies organize their expenses differently. Some have formal RTB/CTB splits.
Others use department-based views. Some just have GL line items.

Your job is to detect:
1. What hierarchy levels exist (e.g., Division > Department > Cost Center > GL Account)
2. Whether there's an explicit Run-the-Business vs Change-the-Business split
3. What expense categories are used and at what level
4. Whether investments/projects are tracked separately
5. What metrics or KPIs are reported alongside expenses

Be thorough but don't hallucinate categories that aren't in the document.

Return ONLY a valid JSON object with this structure:
{
  "structure_tier": <1, 2, or 3 — 1=highly structured, 2=semi-structured, 3=unstructured>,
  "hierarchy_levels": [
    {"name": "<level name>", "depth": <int>, "example_values": ["<val1>", "<val2>"]}
  ],
  "has_rtb_ctb_split": <bool>,
  "has_department_breakdown": <bool>,
  "has_investment_tracking": <bool>,
  "has_kpis": <bool>,
  "detected_categories": [
    {
      "name": "<category name>",
      "level": "<which hierarchy level>",
      "parent": "<parent category name or null>",
      "example_items": ["<item1>", "<item2>"],
      "total_amount": <number or null>
    }
  ],
  "currency": "<detected currency code, e.g. USD>",
  "fiscal_period": "<detected period, e.g. Q3 2024 or FY2024 or null>",
  "confidence": <float 0-1>
}"""

STRUCTURE_DETECTION_USER = """Here is the extracted content from a financial document:

{extracted_text}

Tables found:
{tables_json}

Detect the reporting structure. Return only the JSON object."""

TAXONOMY_MAPPING_SYSTEM = (
    TAXONOMY_REFERENCE
    + """

Given the detected reporting structure below, create a mapping from their categories to our standard SentioCap taxonomy.
For each of their categories:
1. Determine if it's RTB or CTB
2. Map to the most appropriate L2 category
3. Map to L3 domain and L4 activity
4. Rate your confidence (0-1)
5. Flag anything that needs human review (confidence < 0.7)

If a source category spans multiple SentioCap categories, create multiple mappings with allocation_pct (fractions that sum to 1.0).

Return ONLY a valid JSON object:
{
  "mappings": [
    {
      "source_name": "<their category name>",
      "source_level": "<their hierarchy level>",
      "mapped_l1": "RTB" | "CTB",
      "mapped_l2": "<L2 code>",
      "mapped_l3": "<L3 code>",
      "mapped_l4": "<L4 code>",
      "confidence": <float 0-1>,
      "reasoning": "<1-2 sentences>",
      "needs_review": <bool>,
      "allocation_pct": <float, 1.0 if not split>
    }
  ],
  "unmapped": ["<category names that couldn't be mapped>"],
  "summary": "<2-3 sentence summary of the overall mapping>"
}"""
)

TAXONOMY_MAPPING_USER = """Detected reporting structure:
{structure_json}

Map all detected categories to the SentioCap taxonomy. Return only the JSON object."""

FINANCIAL_EXTRACTION_SYSTEM = (
    TAXONOMY_REFERENCE
    + """

You will receive raw financial document content plus a confirmed taxonomy mapping.
Extract actual dollar amounts from the content, apply the taxonomy mapping, and produce classified line items.

Return ONLY a valid JSON object:
{
  "line_items": [
    {
      "source_description": "<original description>",
      "source_category": "<their original category>",
      "amount": <number>,
      "period": "<quarter/year or null>",
      "classified_l1": "RTB" | "CTB",
      "classified_l2": "<L2 code>",
      "classified_l3": "<L3 code>",
      "classified_l4": "<L4 code>",
      "confidence": <float 0-1>
    }
  ],
  "totals": {
    "rtb_total": <number>,
    "ctb_total": <number>,
    "grand_total": <number>,
    "by_l2": {"<L2 code>": <number>}
  },
  "insights": ["<observation 1>", "<observation 2>"]
}"""
)

FINANCIAL_EXTRACTION_USER = """Raw document content:
{extracted_text}

Taxonomy mapping confirmed:
{mapping_json}

Extract all financial line items and apply the mapping. Return only the JSON object."""


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _get_client() -> anthropic.Anthropic:
    return anthropic.Anthropic(api_key=settings.anthropic_api_key)


def _safe_parse(raw: str) -> Any:
    """Strip markdown fences and parse JSON."""
    text = raw.strip()
    if text.startswith("```"):
        parts = text.split("```")
        text = parts[1] if len(parts) > 1 else text
        if text.startswith("json"):
            text = text[4:]
    return json.loads(text.strip())


def _truncate(text: str, max_chars: int = 40_000) -> str:
    """Truncate text for prompt safety."""
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + "\n\n[... truncated for length ...]"


def _detect_numeric_rows(df: pd.DataFrame) -> int:
    """Return index of first row that looks like a header (all-string) before numeric rows."""
    for i, row in df.iterrows():
        str_count = sum(1 for v in row if isinstance(v, str) and v.strip())
        num_count = sum(1 for v in row if isinstance(v, (int, float)) and not pd.isna(v))
        if num_count > str_count:
            return int(i)
    return 0


# ---------------------------------------------------------------------------
# File parsers
# ---------------------------------------------------------------------------


def _parse_pdf(file_path: str) -> tuple[str, list[ExtractedTable], list[Section]]:
    import pdfplumber  # type: ignore

    raw_text_parts: list[str] = []
    tables: list[ExtractedTable] = []
    sections: list[Section] = []

    with pdfplumber.open(file_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            raw_text_parts.append(f"--- Page {page_num} ---\n{text}")

            # Detect sections from headings (lines in ALL CAPS or followed by colon)
            for line in text.splitlines():
                stripped = line.strip()
                if stripped and (stripped.isupper() or stripped.endswith(":")):
                    sections.append(Section(heading=stripped, content="", level=1))

            # Extract tables
            for tbl in page.extract_tables() or []:
                if not tbl or not tbl[0]:
                    continue
                headers = [str(h or "").strip() for h in tbl[0]]
                rows = [[str(c or "").strip() for c in row] for row in tbl[1:]]
                tables.append(
                    ExtractedTable(
                        headers=headers,
                        rows=rows,
                        page_or_slide=page_num,
                    )
                )

    return "\n\n".join(raw_text_parts), tables, sections


def _parse_pptx(file_path: str) -> tuple[str, list[ExtractedTable], list[Section]]:
    from pptx import Presentation  # type: ignore
    from pptx.util import Pt

    prs = Presentation(file_path)
    raw_parts: list[str] = []
    tables: list[ExtractedTable] = []
    sections: list[Section] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        slide_title = ""

        for shape in slide.shapes:
            # Title shapes
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    # Heuristic: large font or placeholder title = heading
                    try:
                        font_size = para.runs[0].font.size
                        is_large = font_size and font_size >= Pt(18)
                    except Exception:
                        is_large = False

                    if shape.shape_type == 13 or is_large or not slide_title:
                        slide_title = slide_title or text
                        sections.append(Section(heading=text, content="", level=1))
                    else:
                        slide_texts.append(text)

            # Table shapes
            if shape.has_table:
                tbl = shape.table
                headers = [tbl.cell(0, c).text.strip() for c in range(tbl.columns.__len__())]
                rows = []
                for r in range(1, tbl.rows.__len__()):
                    rows.append([tbl.cell(r, c).text.strip() for c in range(tbl.columns.__len__())])
                tables.append(
                    ExtractedTable(
                        headers=headers,
                        rows=rows,
                        page_or_slide=slide_num,
                        title=slide_title,
                    )
                )

        raw_parts.append(
            f"--- Slide {slide_num}: {slide_title} ---\n" + "\n".join(slide_texts)
        )

    return "\n\n".join(raw_parts), tables, sections


def _parse_xlsx(file_path: str) -> tuple[str, list[ExtractedTable], list[Section]]:
    raw_parts: list[str] = []
    tables: list[ExtractedTable] = []

    xl = pd.ExcelFile(file_path)
    for sheet_name in xl.sheet_names:
        df = xl.parse(sheet_name, header=None)
        if df.empty:
            continue

        # Find where headers likely start
        header_row = _detect_numeric_rows(df)
        df_data = xl.parse(sheet_name, header=header_row)
        df_data.columns = [str(c).strip() for c in df_data.columns]
        df_data = df_data.dropna(how="all")

        raw_parts.append(f"=== Sheet: {sheet_name} ===\n{df_data.to_string(index=False)}")
        tables.append(
            ExtractedTable(
                headers=list(df_data.columns),
                rows=df_data.values.tolist(),
                page_or_slide=0,
                title=sheet_name,
            )
        )

    return "\n\n".join(raw_parts), tables, []


def _parse_csv(file_path: str) -> tuple[str, list[ExtractedTable], list[Section]]:
    df = pd.read_csv(file_path, dtype=str)
    df = df.fillna("").apply(lambda col: col.str.strip())

    raw_text = df.to_string(index=False)
    tables = [
        ExtractedTable(
            headers=list(df.columns),
            rows=df.values.tolist(),
            page_or_slide=1,
            title=Path(file_path).stem,
        )
    ]
    return raw_text, tables, []


# ---------------------------------------------------------------------------
# DocumentParser
# ---------------------------------------------------------------------------


class DocumentParser:
    """Generalized document intelligence — works with any FP&A format."""

    SUPPORTED_TYPES = {"pdf", "pptx", "xlsx", "xls", "csv"}

    # ----------------------------------------------------------------
    # Step 1: Parse raw content
    # ----------------------------------------------------------------

    async def parse_document(self, file_path: str, file_type: str) -> ParsedDocument:
        """Extract raw text, tables, and sections from the document."""
        ft = file_type.lower().lstrip(".")
        if ft not in self.SUPPORTED_TYPES:
            raise ValueError(f"Unsupported file type: {file_type}")

        try:
            if ft == "pdf":
                raw_text, tables, sections = _parse_pdf(file_path)
                page_count = len(tables)  # rough proxy
            elif ft == "pptx":
                raw_text, tables, sections = _parse_pptx(file_path)
                page_count = len(tables)
            elif ft in ("xlsx", "xls"):
                raw_text, tables, sections = _parse_xlsx(file_path)
                page_count = len(tables)
            else:  # csv
                raw_text, tables, sections = _parse_csv(file_path)
                page_count = 1

            # Detect language heuristically (very light — just ASCII vs non-ASCII ratio)
            non_ascii = sum(1 for c in raw_text if ord(c) > 127)
            lang = "en" if (non_ascii / max(len(raw_text), 1)) < 0.1 else "unknown"

            return ParsedDocument(
                raw_text=raw_text,
                tables=tables,
                sections=sections,
                metadata={
                    "file_type": ft,
                    "page_count": page_count,
                    "table_count": len(tables),
                    "char_count": len(raw_text),
                    "detected_language": lang,
                },
            )
        except Exception as e:
            logger.error(f"parse_document failed for {file_path}: {e}")
            raise

    # ----------------------------------------------------------------
    # Step 2: Detect structure via AI
    # ----------------------------------------------------------------

    async def detect_structure(self, parsed: ParsedDocument) -> DetectedStructure:
        """Use Claude to detect the firm's reporting structure."""
        client = _get_client()

        tables_json = json.dumps(
            [t.to_dict() for t in parsed.tables[:10]],  # limit to first 10
            indent=2,
        )
        extracted_text = _truncate(parsed.raw_text, 30_000)

        user_msg = STRUCTURE_DETECTION_USER.format(
            extracted_text=extracted_text,
            tables_json=tables_json,
        )

        try:
            response = client.messages.create(
                model=MODEL,
                max_tokens=4096,
                system=STRUCTURE_DETECTION_SYSTEM,
                messages=[{"role": "user", "content": user_msg}],
            )
            data = _safe_parse(response.content[0].text)
        except Exception as e:
            logger.error(f"detect_structure AI call failed: {e}")
            data = _fallback_structure_dict()

        # Hydrate dataclasses
        hierarchy_levels = [
            HierarchyLevel(
                name=h.get("name", "Unknown"),
                depth=h.get("depth", 1),
                example_values=h.get("example_values", []),
            )
            for h in data.get("hierarchy_levels", [])
        ]
        detected_categories = [
            DetectedCategory(
                name=c.get("name", ""),
                level=c.get("level", ""),
                parent=c.get("parent"),
                example_items=c.get("example_items", []),
                total_amount=c.get("total_amount"),
            )
            for c in data.get("detected_categories", [])
        ]

        return DetectedStructure(
            structure_tier=data.get("structure_tier", 3),
            hierarchy_levels=hierarchy_levels,
            has_rtb_ctb_split=data.get("has_rtb_ctb_split", False),
            has_department_breakdown=data.get("has_department_breakdown", False),
            has_investment_tracking=data.get("has_investment_tracking", False),
            has_kpis=data.get("has_kpis", False),
            detected_categories=detected_categories,
            currency=data.get("currency", "USD"),
            fiscal_period=data.get("fiscal_period") or "",
            confidence=data.get("confidence", 0.5),
        )

    # ----------------------------------------------------------------
    # Step 3: Map to SentioCap taxonomy via AI
    # ----------------------------------------------------------------

    async def map_to_taxonomy(self, structure: DetectedStructure) -> TaxonomyMapping:
        """Use Claude to map the firm's categories to SentioCap L1-L4 taxonomy."""
        client = _get_client()

        structure_dict = {
            "structure_tier": structure.structure_tier,
            "has_rtb_ctb_split": structure.has_rtb_ctb_split,
            "has_department_breakdown": structure.has_department_breakdown,
            "has_investment_tracking": structure.has_investment_tracking,
            "hierarchy_levels": [
                {"name": h.name, "depth": h.depth, "examples": h.example_values}
                for h in structure.hierarchy_levels
            ],
            "detected_categories": [
                {
                    "name": c.name,
                    "level": c.level,
                    "parent": c.parent,
                    "examples": c.example_items,
                    "total_amount": c.total_amount,
                }
                for c in structure.detected_categories
            ],
        }

        user_msg = TAXONOMY_MAPPING_USER.format(
            structure_json=json.dumps(structure_dict, indent=2)
        )

        try:
            response = client.messages.create(
                model=MODEL,
                max_tokens=4096,
                system=TAXONOMY_MAPPING_SYSTEM,
                messages=[{"role": "user", "content": user_msg}],
            )
            data = _safe_parse(response.content[0].text)
        except Exception as e:
            logger.error(f"map_to_taxonomy AI call failed: {e}")
            data = {"mappings": [], "unmapped": [], "summary": "Mapping failed — manual review required."}

        mappings = [
            CategoryMapping(
                source_name=m.get("source_name", ""),
                source_level=m.get("source_level", ""),
                mapped_l1=m.get("mapped_l1", "RTB"),
                mapped_l2=m.get("mapped_l2", "L2-OPS"),
                mapped_l3=m.get("mapped_l3", "L3-SFT"),
                mapped_l4=m.get("mapped_l4", "L4-SFT-001"),
                confidence=m.get("confidence", 0.5),
                reasoning=m.get("reasoning", ""),
                needs_review=m.get("needs_review", False) or m.get("confidence", 1.0) < 0.7,
                allocation_pct=m.get("allocation_pct", 1.0),
            )
            for m in data.get("mappings", [])
        ]

        return TaxonomyMapping(
            mappings=mappings,
            unmapped=data.get("unmapped", []),
            summary=data.get("summary", ""),
        )

    # ----------------------------------------------------------------
    # Step 4: Extract financials and apply mapping
    # ----------------------------------------------------------------

    async def extract_financials(
        self, parsed: ParsedDocument, mapping: TaxonomyMapping
    ) -> ClassifiedData:
        """Extract dollar amounts from the document and classify them via the mapping."""
        client = _get_client()

        mapping_dict = {
            "mappings": [
                {
                    "source_name": m.source_name,
                    "mapped_l1": m.mapped_l1,
                    "mapped_l2": m.mapped_l2,
                    "mapped_l3": m.mapped_l3,
                    "mapped_l4": m.mapped_l4,
                }
                for m in mapping.mappings
            ]
        }

        extracted_text = _truncate(parsed.raw_text, 30_000)
        user_msg = FINANCIAL_EXTRACTION_USER.format(
            extracted_text=extracted_text,
            mapping_json=json.dumps(mapping_dict, indent=2),
        )

        try:
            response = client.messages.create(
                model=MODEL,
                max_tokens=8192,
                system=FINANCIAL_EXTRACTION_SYSTEM,
                messages=[{"role": "user", "content": user_msg}],
            )
            data = _safe_parse(response.content[0].text)
        except Exception as e:
            logger.error(f"extract_financials AI call failed: {e}")
            data = {"line_items": [], "totals": {}, "insights": ["Extraction failed — manual review required."]}

        line_items = [
            ClassifiedLineItem(
                source_description=li.get("source_description", ""),
                source_category=li.get("source_category", ""),
                amount=float(li.get("amount", 0) or 0),
                period=li.get("period") or "",
                classified_l1=li.get("classified_l1", "RTB"),
                classified_l2=li.get("classified_l2", "L2-OPS"),
                classified_l3=li.get("classified_l3", "L3-SFT"),
                classified_l4=li.get("classified_l4", "L4-SFT-001"),
                confidence=float(li.get("confidence", 0.5) or 0.5),
            )
            for li in data.get("line_items", [])
        ]

        return ClassifiedData(
            line_items=line_items,
            totals=data.get("totals", {}),
            insights=data.get("insights", []),
        )


# ---------------------------------------------------------------------------
# Fallbacks
# ---------------------------------------------------------------------------


def _fallback_structure_dict() -> dict:
    return {
        "structure_tier": 3,
        "hierarchy_levels": [],
        "has_rtb_ctb_split": False,
        "has_department_breakdown": False,
        "has_investment_tracking": False,
        "has_kpis": False,
        "detected_categories": [],
        "currency": "USD",
        "fiscal_period": None,
        "confidence": 0.1,
    }


# Singleton
_parser: Optional[DocumentParser] = None


def get_parser() -> DocumentParser:
    global _parser
    if _parser is None:
        _parser = DocumentParser()
    return _parser
